from typing import List

from pptx import Presentation


def extract_text(pptx_path: str) -> str:
    prs = Presentation(pptx_path)
    lines: List[str] = []

    for index, slide in enumerate(prs.slides, start=1):
        lines.append(f"Slide {index}:")
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            text = text.strip()
            if text:
                lines.append(text)

    return "\n".join(lines).strip()
